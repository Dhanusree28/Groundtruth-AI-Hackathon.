"""
Core orchestration logic for the Automated Insight Engine.
"""

from __future__ import annotations

import argparse
import json
import os
import textwrap
import warnings
from dataclasses import asdict
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional, Tuple

import joblib
import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from sklearn.compose import ColumnTransformer
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    accuracy_score,
    classification_report,
    precision_score,
    recall_score,
    roc_auc_score,
)
from sklearn.model_selection import train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import OneHotEncoder, StandardScaler

from pptx import Presentation
from pptx.util import Inches, Pt

from src.config import MODEL_CONFIG, PATHS

sns.set_theme(style="whitegrid")


class InsightEngine:
    """End-to-end automation engine."""

    def __init__(
        self,
        paths=PATHS,
        model_config=MODEL_CONFIG,
        llm_provider: str = "openai",
    ) -> None:
        self.paths = paths
        self.config = model_config
        self.llm_provider = llm_provider
        self._ensure_directories()

    def _ensure_directories(self) -> None:
        for directory in (
            self.paths.data_dir,
            self.paths.reports_dir,
            self.paths.models_dir,
            self.paths.artifacts_dir,
            self.paths.uploads_dir,
        ):
            directory.mkdir(parents=True, exist_ok=True)

    def load_data(self, file_path: Optional[Path] = None) -> pd.DataFrame:
        data_path = Path(file_path) if file_path else self.paths.data_dir / "sample_kag_conversion.csv"
        if not data_path.exists():
            raise FileNotFoundError(
                f"Could not locate dataset at {data_path}. "
                "Place KAG_conversion_data.csv in the data directory or upload via the UI."
            )
        df = pd.read_csv(data_path)
        df.columns = [c.strip().lower() for c in df.columns]

        cols = set(df.columns)

        # Auto-detect schema and adapt configuration
        # 1) Message/label style datasets
        if {"message", "label"}.issubset(cols):
            self.config.target_column = "label"
            self.config.numeric_features = []
            self.config.categorical_features = ["message"]

        # 2) KAG conversion style datasets (default)
        elif {"approved_conversion", "impressions", "clicks"}.issubset(cols):
            self.config.target_column = "approved_conversion"
            self.config.numeric_features = [
                "impressions",
                "clicks",
                "spent",
                "total_conversion",
                "interest",
            ]
            self.config.categorical_features = [
                "age",
                "gender",
                "xyz_campaign_id",
                "fb_campaign_id",
            ]

        target = self.config.target_column
        if target not in df.columns:
            if "converted" in df.columns:
                df[target] = df["converted"]
            elif "approved_conversion" in df.columns:
                df[target] = df["approved_conversion"]
            else:
                raise KeyError(
                    f"Could not find target column '{target}'. Available columns: {list(df.columns)}"
                )

        # Target preparation
        if target == "label":
            labels = df[target].astype(str).str.strip().str.lower()
            mapping = {
                "ham": 0,
                "spam": 1,
                "no": 0,
                "yes": 1,
                "0": 0,
                "1": 1,
                "false": 0,
                "true": 1,
                "negative": 0,
                "positive": 1,
            }
            if labels.isin(mapping).all():
                df[target] = labels.map(mapping).astype(int)
            else:
                uniques = sorted(labels.unique())
                if len(uniques) == 2:
                    auto_map = {uniques[0]: 0, uniques[1]: 1}
                    df[target] = labels.map(auto_map).astype(int)
                else:
                    raise ValueError(
                        f"Could not map label values to 0/1. Unique labels: {uniques}"
                    )
        else:
            df[target] = (pd.to_numeric(df[target], errors="coerce") > 0).astype(int)

        for column in self.config.numeric_features:
            if column in df.columns:
                df[column] = pd.to_numeric(df[column], errors="coerce")
        df = df.dropna(subset=[target])
        return df

    def _build_pipeline(self) -> Pipeline:
        transformers = ColumnTransformer(
            transformers=[
                ("numeric", StandardScaler(), self.config.numeric_features),
                (
                    "categorical",
                    OneHotEncoder(handle_unknown="ignore"),
                    self.config.categorical_features,
                ),
            ]
        )

        return Pipeline(
            steps=[
                ("preprocess", transformers),
                (
                    "model",
                    LogisticRegression(
                        max_iter=1000,
                        class_weight="balanced",
                        random_state=self.config.random_state,
                    ),
                ),
            ]
        )

    def train_model(self, df: pd.DataFrame) -> Tuple[Pipeline, Dict[str, float], pd.DataFrame]:
        X = df[self.config.numeric_features + self.config.categorical_features]
        y = df[self.config.target_column]

        label_counts = y.value_counts()
        if len(label_counts) < 2:
            raise ValueError(
                "Target column contains only one class after preprocessing. "
                "Provide a dataset with both positive and negative examples."
            )

        stratify_labels = y if label_counts.min() >= 2 else None
        if stratify_labels is None:
            warnings.warn(
                "Insufficient samples for one of the classes; falling back to unstratified split.",
                RuntimeWarning,
            )

        X_train, X_test, y_train, y_test = train_test_split(
            X,
            y,
            test_size=self.config.test_size,
            random_state=self.config.random_state,
            stratify=stratify_labels,
        )

        pipeline = self._build_pipeline()
        pipeline.fit(X_train, y_train)
        preds = pipeline.predict(X_test)
        proba = pipeline.predict_proba(X_test)[:, 1]

        metrics = {
            "accuracy": round(float(accuracy_score(y_test, preds)), 3),
            "precision": round(float(precision_score(y_test, preds)), 3),
            "recall": round(float(recall_score(y_test, preds)), 3),
            "roc_auc": round(float(roc_auc_score(y_test, proba)), 3),
        }

        joblib.dump(pipeline, self.paths.models_dir / "conversion_model.joblib")
        (self.paths.models_dir / "classification_report.json").write_text(
            json.dumps(classification_report(y_test, preds, output_dict=True), indent=2)
        )

        return pipeline, metrics, X_test.assign(actual=y_test, predicted=preds)

    def _conversion_summary(self, df: pd.DataFrame) -> Dict[str, float | Dict[str, float]]:
        target = self.config.target_column
        summary: Dict[str, Dict[str, float] | float] = {
            "overall_rate": round(float(df[target].mean()), 3)
        }

        # Generic label distribution (works for message/label datasets)
        if target == "label":
            summary["by_label"] = (
                df[target].value_counts(normalize=True).round(3).to_dict()
            )

        if "age" in df.columns:
            summary["by_age"] = (
                df.groupby("age")[target].mean().sort_values(ascending=False).round(3).to_dict()
            )
        if "gender" in df.columns:
            summary["by_gender"] = (
                df.groupby("gender")[target].mean().sort_values(ascending=False).round(3).to_dict()
            )
        if "xyz_campaign_id" in df.columns:
            campaign_perf = (
                df.groupby("xyz_campaign_id")[target]
                .mean()
                .sort_values(ascending=False)
                .head(5)
                .round(3)
                .to_dict()
            )
            summary["top_campaigns"] = campaign_perf

            spend_eff = (
                df.groupby("xyz_campaign_id")
                .agg(conversions=(target, "sum"), spend=("spent", "sum"))
            )
            spend_eff["conv_per_dollar"] = spend_eff["conversions"] / spend_eff["spend"].replace(0, pd.NA)
            summary["efficient_campaigns"] = (
                spend_eff["conv_per_dollar"].dropna().sort_values(ascending=False).head(5).round(3).to_dict()
            )

        return summary

    def _create_visual(self, df: pd.DataFrame, column: str, title: str, file_name: str) -> Path:
        fig, ax = plt.subplots(figsize=(6, 4))
        target = self.config.target_column
        chart_data = df.groupby(column)[target].mean().sort_values(ascending=False).head(10)
        sns.barplot(x=chart_data.values, y=chart_data.index, palette="viridis", ax=ax)
        ax.set_title(title)
        ax.set_xlabel("Conversion rate")
        ax.set_xlim(0, 1)
        fig.tight_layout()
        output_path = self.paths.artifacts_dir / file_name
        fig.savefig(output_path, dpi=300)
        plt.close(fig)
        return output_path

    def _generate_llm_text(self, metrics: Dict[str, float], summary: Dict[str, float | Dict[str, float]]) -> str:
        if self.llm_provider != "openai":
            return self._fallback_summary(metrics, summary)

        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return self._fallback_summary(metrics, summary)

        try:
            from openai import OpenAI

            client = OpenAI(api_key=api_key)
            prompt = textwrap.dedent(
                f"""
                You are a marketing analytics assistant. Craft a crisp executive summary (<=150 words)
                using the following metrics: {json.dumps(metrics)} and cohort summary: {json.dumps(summary)}.
                Highlight what is working, what is not, and one next action.
                """
            )
            response = client.responses.create(
                model="gpt-4o-mini",
                input=prompt,
                max_output_tokens=220,
            )
            return response.output[0].content[0].text or self._fallback_summary(metrics, summary)
        except Exception:
            return self._fallback_summary(metrics, summary)

    def _fallback_summary(self, metrics: Dict[str, float], summary: Dict[str, float | Dict[str, float]]) -> str:
        def pick(bucket: Optional[Dict[str, float]], reverse: bool = True) -> Tuple[str, float]:
            if not isinstance(bucket, dict) or not bucket:
                return ("N/A", 0.0)
            key = max(bucket, key=bucket.get) if reverse else min(bucket, key=bucket.get)
            return key, bucket[key]

        best_age, best_age_val = pick(summary.get("by_age"))
        best_gender, _ = pick(summary.get("by_gender"))
        best_campaign, best_campaign_val = pick(summary.get("top_campaigns"))
        weak_campaign, weak_campaign_val = pick(summary.get("top_campaigns"), reverse=False)

        return textwrap.dedent(
            f"""
            Overall approval conversion rate is {summary['overall_rate']:.1%} with model precision {metrics['precision']:.1%}
            and recall {metrics['recall']:.1%}. {best_age} audiences respond best, followed by {best_gender} users.
            Campaign {best_campaign} is outperforming at {best_campaign_val:.1%}, while campaign {weak_campaign} trails at {weak_campaign_val:.1%}.
            Reinvest into the winning creative mix and pressure-test underperforming campaigns with fresh copy and tighter targeting.
            """
        ).strip()

    def _build_pdf(
        self,
        metrics: Dict[str, float],
        summary: Dict[str, float | Dict[str, float]],
        narrative: str,
        visuals: Dict[str, Path],
    ) -> Path:
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        pdf_path = self.paths.reports_dir / f"insight_report_{timestamp}.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=LETTER)
        width, height = LETTER
        margin = 0.75 * inch

        c.setFont("Helvetica-Bold", 18)
        c.drawString(margin, height - margin, "Automated Insight Report")
        c.setFont("Helvetica", 10)
        c.drawString(margin, height - margin - 14, f"Generated: {datetime.utcnow().isoformat()} UTC")

        y = height - margin - 40
        c.setFont("Helvetica-Bold", 12)
        c.drawString(margin, y, "Model KPIs")
        c.setFont("Helvetica", 10)
        y -= 14
        for metric, value in metrics.items():
            c.drawString(margin, y, f"{metric.title()}: {value:.3f}")
            y -= 12

        y -= 10
        c.setFont("Helvetica-Bold", 12)
        c.drawString(margin, y, "Narrative Summary")
        y -= 14
        text_obj = c.beginText(margin, y)
        text_obj.setFont("Helvetica", 10)
        for line in textwrap.wrap(narrative, width=90):
            text_obj.textLine(line)
        c.drawText(text_obj)

        current_y = y - 80
        for title, path in visuals.items():
            if current_y < 1.5 * inch:
                c.showPage()
                current_y = height - margin
            c.setFont("Helvetica-Bold", 12)
            c.drawString(margin, current_y, title)
            current_y -= 12
            c.drawImage(str(path), margin, current_y - 3.0 * inch, width=width - 2 * margin, height=3.0 * inch)
            current_y -= 3.2 * inch

        c.showPage()
        c.setFont("Helvetica-Bold", 12)
        c.drawString(margin, height - margin - 20, "Appendix: Cohort Summary")
        appendix_y = height - margin - 40
        for bucket, values in summary.items():
            c.setFont("Helvetica-Bold", 11)
            c.drawString(margin, appendix_y, bucket.replace("_", " ").title())
            appendix_y -= 12
            c.setFont("Helvetica", 10)
            if isinstance(values, dict):
                for key, val in values.items():
                    c.drawString(margin + 12, appendix_y, f"{key}: {val:.1%}")
                    appendix_y -= 12
            else:
                c.drawString(margin + 12, appendix_y, f"{values:.1%}")
                appendix_y -= 12
            appendix_y -= 6
        c.save()
        return pdf_path

    def _build_ppt(
        self,
        metrics: Dict[str, float],
        summary: Dict[str, float | Dict[str, float]],
        narrative: str,
        visuals: Dict[str, Path],
    ) -> Path:
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        ppt_path = self.paths.reports_dir / f"insight_report_{timestamp}.pptx"
        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Automated Insight Report"
        slide.placeholders[1].text = f"Generated {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"

        kpi_slide = prs.slides.add_slide(prs.slide_layouts[1])
        kpi_slide.shapes.title.text = "Model KPIs"
        text_frame = kpi_slide.shapes.placeholders[1].text_frame
        text_frame.clear()
        for metric, value in metrics.items():
            p = text_frame.add_paragraph()
            p.text = f"{metric.title()}: {value:.3f}"
            p.level = 0

        narrative_slide = prs.slides.add_slide(prs.slide_layouts[1])
        narrative_slide.shapes.title.text = "Narrative Insights"
        body = narrative_slide.shapes.placeholders[1].text_frame
        body.text = ""
        for chunk in textwrap.wrap(narrative, width=70):
            p = body.add_paragraph()
            p.text = chunk
            p.level = 0

        for title, path in visuals.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            slide.shapes.add_picture(str(path), Inches(1), Inches(1.5), width=Inches(8))

        appendix_slide = prs.slides.add_slide(prs.slide_layouts[1])
        appendix_slide.shapes.title.text = "Cohort Summary"
        appendix_frame = appendix_slide.shapes.placeholders[1].text_frame
        appendix_frame.text = ""
        for bucket, values in summary.items():
            p = appendix_frame.add_paragraph()
            p.text = bucket.replace("_", " ").title()
            p.level = 0
            if isinstance(values, dict):
                for key, val in values.items():
                    entry = appendix_frame.add_paragraph()
                    entry.text = f"{key}: {val:.1%}"
                    entry.level = 1
        prs.save(ppt_path)
        return ppt_path

    def build_report(
        self,
        metrics: Dict[str, float],
        summary: Dict[str, Dict[str, float]],
        narrative: str,
        visuals: Dict[str, Path],
        report_type: str = "PDF",
    ) -> Path:
        if report_type.upper() == "PDF":
            return self._build_pdf(metrics, summary, narrative, visuals)
        if report_type.upper() in {"PPT", "PPTX", "POWERPOINT"}:
            return self._build_ppt(metrics, summary, narrative, visuals)
        raise ValueError("report_type must be either 'PDF' or 'PPTX'")

    def run(self, file_path: Optional[Path] = None, report_type: str = "PDF") -> Dict[str, Path]:
        df = self.load_data(file_path)
        pipeline, metrics, eval_sample = self.train_model(df)
        summary = self._conversion_summary(df)

        visuals = {}
        if "age" in df.columns:
            visuals["Conversion by Age"] = self._create_visual(df, "age", "Conversion by Age", "age.png")
        if "gender" in df.columns:
            visuals["Conversion by Gender"] = self._create_visual(df, "gender", "Conversion by Gender", "gender.png")
        if "xyz_campaign_id" in df.columns:
            visuals["Top Campaigns"] = self._create_visual(
                df, "xyz_campaign_id", "Top Campaign Conversion Rate", "campaign.png"
            )

        narrative = self._generate_llm_text(metrics, summary)
        report_path = self.build_report(metrics, summary, narrative, visuals, report_type=report_type)

        eval_path = self.paths.artifacts_dir / "eval_sample.csv"
        eval_sample.to_csv(eval_path, index=False)

        return {
            "report": report_path,
            "model": self.paths.models_dir / "conversion_model.joblib",
            "evaluation_sample": eval_path,
            "metrics": metrics,
            "narrative": narrative,
        }


def run_from_cli(
    data_path: Optional[str] = None,
    report_type: str = "PDF",
) -> Dict[str, Path]:
    engine = InsightEngine()
    return engine.run(file_path=Path(data_path) if data_path else None, report_type=report_type)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Automated Insight Engine CLI")
    parser.add_argument("--data_path", type=str, default=None, help="Path to CSV dataset")
    parser.add_argument(
        "--report_type",
        type=str,
        default="PDF",
        choices=["PDF", "PPTX", "PPT", "PowerPoint"],
        help="Output format",
    )
    args = parser.parse_args()

    artifacts = run_from_cli(data_path=args.data_path, report_type=args.report_type)
    print(json.dumps({k: str(v) if isinstance(v, Path) else v for k, v in artifacts.items()}, indent=2))

